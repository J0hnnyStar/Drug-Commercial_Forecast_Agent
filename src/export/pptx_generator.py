"""
PPTX One-Pager Generator for Pharmaceutical Commercial Forecasting.

Creates professional PowerPoint slide with key insights and recommendations.
Following Linus principle: Simple, direct, single source of truth.
"""
import io
from typing import Dict, Any, Optional
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

def create_commercial_forecast_slide(
    analysis_result: Dict[str, Any],
    template_path: Optional[str] = None
) -> io.BytesIO:
    """
    Create professional PPTX one-pager from analysis results.
    
    Args:
        analysis_result: Complete analysis results from AI system
        template_path: Optional PowerPoint template path
    
    Returns:
        io.BytesIO: PPTX file as bytes buffer
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")
    
    # Create presentation
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    
    # Add title slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Extract data
    query = analysis_result.get("query", "Pharmaceutical Investment Analysis")
    recommendation = analysis_result.get("recommendation", {})
    parameters = analysis_result.get("parameters", {})
    characteristics = analysis_result.get("characteristics", {})
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Pharmaceutical Commercial Forecast"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)  # Navy blue
    
    # Subtitle with query
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f'Query: "{query}"'
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(14)
    subtitle_para.font.italic = True
    
    # Recommendation section (left column)
    rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
    rec_frame = rec_box.text_frame
    
    decision = recommendation.get("decision", "No Decision")
    rationale = recommendation.get("rationale", "No rationale available")
    confidence = recommendation.get("confidence", "unknown")
    
    rec_frame.text = "INVESTMENT RECOMMENDATION"
    rec_para = rec_frame.paragraphs[0]
    rec_para.font.size = Pt(16)
    rec_para.font.bold = True
    rec_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # Decision with color coding
    decision_para = rec_frame.add_paragraph()
    decision_para.text = f"Decision: {decision}"
    decision_para.font.size = Pt(14)
    decision_para.font.bold = True
    if decision in ["STRONG GO", "GO"]:
        decision_para.font.color.rgb = RGBColor(0, 128, 0)  # Green
    elif decision == "CONDITIONAL GO":
        decision_para.font.color.rgb = RGBColor(255, 140, 0)  # Orange
    else:
        decision_para.font.color.rgb = RGBColor(220, 20, 60)  # Red
    
    # Confidence
    conf_para = rec_frame.add_paragraph()
    conf_para.text = f"Confidence: {confidence.upper()}"
    conf_para.font.size = Pt(12)
    
    # Rationale
    rationale_para = rec_frame.add_paragraph()
    rationale_para.text = f"Rationale: {rationale}"
    rationale_para.font.size = Pt(11)
    
    # Key Metrics section (right column)
    metrics_box = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(4))
    metrics_frame = metrics_box.text_frame
    
    metrics_frame.text = "KEY FINANCIAL METRICS"
    metrics_para = metrics_frame.paragraphs[0]
    metrics_para.font.size = Pt(16)
    metrics_para.font.bold = True
    metrics_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # Extract metrics
    key_metrics = recommendation.get("key_metrics", {})
    npv = key_metrics.get("npv_billions", 0)
    success_rate = key_metrics.get("success_rate", 0)
    market_size = key_metrics.get("market_size", 0)
    
    # NPV
    npv_para = metrics_frame.add_paragraph()
    npv_para.text = f"NPV: ${npv:.1f}B"
    npv_para.font.size = Pt(14)
    npv_para.font.bold = True
    npv_para.font.color.rgb = RGBColor(0, 128, 0) if npv > 0 else RGBColor(220, 20, 60)
    
    # Success Rate
    success_para = metrics_frame.add_paragraph()
    success_para.text = f"Success Rate: {success_rate:.0%}"
    success_para.font.size = Pt(12)
    
    # Market Size
    market_para = metrics_frame.add_paragraph()
    market_para.text = f"Market Size: {market_size:,} patients"
    market_para.font.size = Pt(12)
    
    # Pricing info if available
    pricing = parameters.get("pricing", {})
    if pricing:
        pricing_para = metrics_frame.add_paragraph()
        pricing_para.text = f"Price: ${pricing.get('list_price', 0):,}/month"
        pricing_para.font.size = Pt(12)
        
        access_para = metrics_frame.add_paragraph()
        access_para.text = f"Access: {pricing.get('access_tier', 'Unknown')}"
        access_para.font.size = Pt(12)
    
    # Drug characteristics section (bottom)
    char_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
    char_frame = char_box.text_frame
    
    char_frame.text = "DRUG PROFILE"
    char_para = char_frame.paragraphs[0]
    char_para.font.size = Pt(14)
    char_para.font.bold = True
    char_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # Create profile string
    profile_items = []
    if characteristics.get("drug_type"):
        profile_items.append(f"Type: {characteristics['drug_type']}")
    if characteristics.get("indication_area"):
        profile_items.append(f"Area: {characteristics['indication_area']}")
    if characteristics.get("patient_population"):
        profile_items.append(f"Population: {characteristics['patient_population']}")
    if characteristics.get("competitive_position"):
        profile_items.append(f"Position: {characteristics['competitive_position']}")
    
    profile_para = char_frame.add_paragraph()
    profile_para.text = " • ".join(profile_items) if profile_items else "No drug profile available"
    profile_para.font.size = Pt(11)
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = f"Generated by AI Commercial Forecast Agent • {datetime.now().strftime('%B %d, %Y')}"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(9)
    footer_para.font.color.rgb = RGBColor(128, 128, 128)
    
    # Save to bytes buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return buffer

def generate_pptx_filename(query: str) -> str:
    """Generate clean filename from query."""
    # Clean query for filename
    clean_query = "".join(c for c in query if c.isalnum() or c in " -_").strip()
    clean_query = clean_query.replace(" ", "_")[:50]  # Limit length
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    return f"commercial_forecast_{clean_query}_{timestamp}.pptx"

if __name__ == "__main__":
    # Test with sample data
    sample_result = {
        "query": "Should we develop a Tezspire competitor for pediatric severe asthma?",
        "recommendation": {
            "decision": "CONDITIONAL GO",
            "rationale": "Strong market need but high development risk",
            "confidence": "medium",
            "key_metrics": {
                "npv_billions": 1.8,
                "success_rate": 0.65,
                "market_size": 85000
            }
        },
        "parameters": {
            "pricing": {
                "list_price": 4200,
                "access_tier": "PA",
                "gtn_pct": 0.72
            }
        },
        "characteristics": {
            "drug_type": "Biologic",
            "indication_area": "Respiratory",
            "patient_population": "Pediatric",
            "competitive_position": "Me-too"
        }
    }
    
    try:
        buffer = create_commercial_forecast_slide(sample_result)
        print(f"Generated PPTX: {len(buffer.getvalue())} bytes")
        
        # Save test file
        with open("test_forecast.pptx", "wb") as f:
            f.write(buffer.getvalue())
        print("Saved test_forecast.pptx")
        
    except Exception as e:
        print(f"Test failed: {e}")